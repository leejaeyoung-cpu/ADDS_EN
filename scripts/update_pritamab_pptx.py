"""
Update pritamab ADDS PPTX with newly generated figures
- Replace all slide images with academic-style English figures
- Update Supplementary Table S1 parameters (w1, w2, w3, k)
- Remove Fig 5-A panel (keep B only)
- Fix text boxes (Korean → English)
"""

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import copy
import os

SRC_PPTX  = r"f:\ADDS\26.02.21 pritamab ADDS 결과 정리.pptx"
DST_PPTX  = r"f:\ADDS\26.02.22 pritamab ADDS 결과 정리_revised.pptx"
FIG_DIR   = r"f:\ADDS\outputs\pritamab_pptx_figures"

# Map: slide index (0-based) → (new image path, label, position mode)
# position mode: 'fill' = full slide, 'center' = centered with margins
SLIDE_UPDATES = {
    0: {'img': 'fig1_AB.png',       'label': 'Figure 1'},  # Slide 1
    1: {'img': 'fig2_AB.png',       'label': 'Figure 2'},  # Slide 2
    2: {'img': 'fig3_AB.png',       'label': 'Figure 3'},  # Slide 3
    3: {'img': 'fig4.png',          'label': 'Figure 4'},  # Slide 4
    4: {'img': 'fig5_B.png',        'label': 'Figure 5'},  # Slide 5 (B only)
    5: {'img': 'fig6.png',          'label': 'Figure 6'},  # Slide 6
    6: None,                                                 # Slide 7: Table only
    7: {'img': 'fig7_summary.png',  'label': 'Figure 7'},  # Slide 8
}

# Supplementary Table S1 parameter values to fill in
TABLE_FILLS = {
    # (row_index, col_index): new_value
    (1, 3): '0.50',   # w1 fixed value
    (2, 3): '0.30',   # w2 fixed value
    (3, 3): '0.20',   # w3 fixed value
    (4, 3): '10',     # T_max (already there, but confirm)
    (6, 3): '1.50',   # k fixed value
    # "Score formula" – update note
    (1, 4): 'E_pred scaled to 0–1',
    (2, 4): 'S_pred scaled to 0–1',
    (3, 4): 'Penalty uses normalized toxicity',
    (4, 4): 'T_tox scaled to 1–10',
    (6, 4): 'Larger k → larger HR shift per ΔScore',
}


def replace_slide_image(slide, new_img_path, label_text=None):
    """Remove all existing pictures and add new full-slide image."""
    slide_w = slide.shapes[0].left if slide.shapes else Inches(13.33)

    # Get slide dimensions from presentation
    # We'll use a standard ratio based on the shapes present
    # Remove existing picture shapes
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            shapes_to_remove.append(shape)
        # Remove highlight rectangles (AutoShapes with no text)
        elif shape.shape_type == 1:  # AUTO_SHAPE
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                shapes_to_remove.append(shape)

    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Update label text box (Figure N)
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt.startswith('Figure') and label_text:
                shape.text_frame.paragraphs[0].runs[0].text = label_text
                for para in shape.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.LEFT
                    for run in para.runs:
                        run.font.bold = True
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(31, 73, 125)  # dark blue

    # Add new picture – fill slide with margin
    # Use fixed slide dimensions (widescreen 13.33" × 7.5")
    sw = Inches(13.33)
    sh = Inches(7.5)
    margin = Inches(0.3)
    label_h = Inches(0.5)

    pic = slide.shapes.add_picture(
        new_img_path,
        left=margin,
        top=label_h + margin,
        width=sw - 2 * margin,
        height=sh - label_h - 2 * margin
    )
    return pic


def fill_table_parameters(slide):
    """Fill in blank parameter values in Supplementary Table S1."""
    for shape in slide.shapes:
        if shape.shape_type == 19:  # TABLE
            table = shape.table
            for (r, c), val in TABLE_FILLS.items():
                try:
                    cell = table.cell(r, c)
                    current = cell.text.strip()
                    if current in ('', '[ ]', '[ ]') or 'scaled' in current or 'Larger' in current:
                        # Clear existing text and set new value
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ''
                        # Set value in first paragraph
                        tf = cell.text_frame
                        if tf.paragraphs:
                            p = tf.paragraphs[0]
                            if p.runs:
                                p.runs[0].text = val
                            else:
                                from pptx.util import Pt
                                run = p.add_run()
                                run.text = val
                                run.font.size = Pt(9)
                        print(f"    Table [{r},{c}] → {val}")
                except Exception as e:
                    print(f"    Table [{r},{c}] error: {e}")


def fix_text_boxes(slide, label_text):
    """Ensure all text boxes use English labels."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    # Fix common Korean labels
                    replacements = {
                        '평균 독성 점수': 'Mean Toxicity Score',
                        '독성 점수': 'Toxicity Score',
                        '관리 가능한 독성': 'Manageable Toxicity',
                        '임상유효 기준': 'Clinical Efficacy Threshold',
                        '임상': 'Clinical',
                        '생존율': 'Survival Rate',
                        '위험비': 'Hazard Ratio',
                        '대조군': 'Control Group',
                        '실험군': 'Treatment Group',
                    }
                    for kor, eng in replacements.items():
                        if kor in t:
                            run.text = t.replace(kor, eng)
                            print(f"    Text: '{kor}' → '{eng}'")


def main():
    print("Loading PPTX...")
    prs = Presentation(SRC_PPTX)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"Slide size: {slide_w.inches:.2f}\" × {slide_h.inches:.2f}\"")

    for slide_idx, config in SLIDE_UPDATES.items():
        slide = prs.slides[slide_idx]
        print(f"\n[Slide {slide_idx+1}]", end=' ')

        if config is None:
            # Slide 7: Table only – fill parameters
            print("Supplementary Table S1 – filling parameters...")
            fill_table_parameters(slide)
        else:
            img_path = os.path.join(FIG_DIR, config['img'])
            label    = config['label']
            print(f"{label} → {config['img']}")

            if not os.path.exists(img_path):
                print(f"  WARNING: Image not found: {img_path}")
                continue

            # Fix text boxes first
            fix_text_boxes(slide, label)

            # Replace image
            replace_slide_image(slide, img_path, label_text=label)
            print(f"  Image replaced")

    prs.save(DST_PPTX)
    print(f"\n=== Saved: {DST_PPTX} ===")


if __name__ == '__main__':
    main()
