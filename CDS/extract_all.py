import os
import io
import sys
import easyocr
import warnings
from pptx import Presentation

warnings.filterwarnings('ignore')

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

def extract_text_from_image(image_path):
    try:
        reader = easyocr.Reader(['ko', 'en'], verbose=False)
        result = reader.readtext(image_path, detail=0)
        return "\n".join(result)
    except Exception as e:
        return f"Error extracting from {image_path}: {e}"

def extract_text_from_shape(shape):
    text = ""
    if hasattr(shape, "text"):
        text += shape.text + "\n"
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            text += extract_text_from_shape(sub_shape)
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                text += cell.text + " "
            text += "\n"
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text_runs.append(extract_text_from_shape(shape))
        return "\n".join(text_runs)
    except Exception as e:
        return f"Error extracting from {pptx_path}: {e}"

def main():
    folder = "."
    output_file = "all_extracted_text.txt"
    with open(output_file, "w", encoding="utf-8") as f:
        for file in os.listdir(folder):
            path = os.path.join(folder, file)
            if file.endswith(".png"):
                text = extract_text_from_image(path)
                f.write(f"--- {file} ---\n{text}\n\n")
            elif file.endswith(".pptx") and not file.startswith("~"):
                text = extract_text_from_pptx(path)
                f.write(f"--- {file} ---\n{text}\n\n")

if __name__ == "__main__":
    main()
