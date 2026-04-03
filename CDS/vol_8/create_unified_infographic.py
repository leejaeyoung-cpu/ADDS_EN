import os
import math
from PIL import Image, ImageDraw, ImageFont, ImageFilter

W, H = 1920, 1080
OUT_PATH = r"f:\ADDS\CDS\ADDS_Unified_Architecture_Slide.png"
PPT_PATH = r"f:\ADDS\CDS\ADDS_Unified_Architecture_Slide.pptx"

# Source Images
IMG1 = r"C:\Users\brook\.gemini\antigravity\brain\574d389b-b6ed-48ef-afb4-dd86c25c56c1\media__1774072909206.png" # Cellpose
IMG2 = r"C:\Users\brook\.gemini\antigravity\brain\574d389b-b6ed-48ef-afb4-dd86c25c56c1\media__1774072909234.jpg" # System

def get_font(size, bold=False):
    font_name = "malgunbd.ttf" if bold else "malgun.ttf"
    font_path = os.path.join(os.environ.get("WINDIR", "C:\\Windows"), "Fonts", font_name)
    try:
        return ImageFont.truetype(font_path, size)
    except:
        return ImageFont.load_default()

def create_gradient_bg(w, h, c1, c2):
    base = Image.new("RGBA", (w, h), c1)
    draw = ImageDraw.Draw(base)
    for y in range(h):
        r = int(c1[0] + (c2[0] - c1[0]) * y / h)
        g = int(c1[1] + (c2[1] - c1[1]) * y / h)
        b = int(c1[2] + (c2[2] - c1[2]) * y / h)
        draw.line([(0, y), (w, y)], fill=(r, g, b, 255))
    return base

def add_drop_shadow(image, offset=(10, 15), blur_radius=25, shadow_color=(0,0,0,180)):
    shadow = Image.new('RGBA', image.size, shadow_color)
    padded_size = (image.width + offset[0] + blur_radius*2,
                   image.height + offset[1] + blur_radius*2)
    base = Image.new('RGBA', padded_size, (0,0,0,0))
    base.paste(shadow, (blur_radius + offset[0], blur_radius + offset[1]))
    base = base.filter(ImageFilter.GaussianBlur(blur_radius))
    base.paste(image, (blur_radius, blur_radius), image if image.mode == 'RGBA' else None)
    return base

def draw_arrow(draw, start, end, fill, width=8, arrow_size=20):
    draw.line([start, end], fill=fill, width=width)
    angle = math.atan2(end[1]-start[1], end[0]-start[0])
    p1 = (end[0] - arrow_size * math.cos(angle - math.pi/6),
          end[1] - arrow_size * math.sin(angle - math.pi/6))
    p2 = (end[0] - arrow_size * math.cos(angle + math.pi/6),
          end[1] - arrow_size * math.sin(angle + math.pi/6))
    draw.polygon([end, p1, p2], fill=fill)

def main():
    print("Generating Unified Infographic with Source Images...")
    # Dark modern background
    canvas = create_gradient_bg(W, H, (10, 15, 29), (26, 43, 76))
    draw = ImageDraw.Draw(canvas, "RGBA")

    # Header
    draw.text((80, 50), "ADDS: Unified AI Decision Support System", font=get_font(52, True), fill=(255, 255, 255))
    draw.text((80, 120), "Seamless integration of high-fidelity cellular segmentation and multi-modal synergy prediction.", font=get_font(26, False), fill=(148, 163, 184))

    # Load Source Images
    try:
        im1 = Image.open(IMG1).convert("RGBA")
        im2 = Image.open(IMG2).convert("RGBA")
    except Exception as e:
        print("Error loading images:", e)
        return

    # Resize Image 1 (Cellpose) - Wide image, let's scale it to fit the left half
    # Max width ~820px, Max height ~600px
    w1, h1 = im1.size
    scale1 = min(840 / w1, 600 / h1)
    nw1, nh1 = int(w1 * scale1), int(h1 * scale1)
    im1_res = im1.resize((nw1, nh1), Image.LANCZOS)
    
    # Resize Image 2 (System Pipeline)
    w2, h2 = im2.size
    scale2 = min(840 / w2, 600 / h2)
    nw2, nh2 = int(w2 * scale2), int(h2 * scale2)
    im2_res = im2.resize((nw2, nh2), Image.LANCZOS)

    # Add drop shadows
    im1_shadowed = add_drop_shadow(im1_res)
    im2_shadowed = add_drop_shadow(im2_res)

    # Card backgrounds (Container for the images)
    card_y = 200
    card_h = 750
    
    # Left Card
    c1_x = 60
    c1_w = 860
    draw.rounded_rectangle([c1_x, card_y, c1_x+c1_w, card_y+card_h], radius=20, fill=(30, 41, 59, 200), outline=(56, 189, 248, 255), width=2)
    draw.text((c1_x + 40, card_y + 30), "Step 1. AI Pathology Engine (Cellpose Processing)", font=get_font(28, True), fill=(56, 189, 248))
    
    # Paste Image 1 centered in its card
    pos1_x = c1_x + (c1_w - im1_shadowed.width) // 2
    pos1_y = card_y + 100 + (card_h - 100 - im1_shadowed.height) // 2
    canvas.paste(im1_shadowed, (pos1_x, pos1_y), im1_shadowed)

    # Right Card
    c2_x = c1_x + c1_w + 80   # 60 + 860 + 80 = 1000
    c2_w = 860
    draw.rounded_rectangle([c2_x, card_y, c2_x+c2_w, card_y+card_h], radius=20, fill=(30, 41, 59, 200), outline=(16, 185, 129, 255), width=2)
    draw.text((c2_x + 40, card_y + 30), "Step 2. Therapy Selection & Integration Pipeline", font=get_font(28, True), fill=(16, 185, 129))
    
    # Paste Image 2 centered in its card
    pos2_x = c2_x + (c2_w - im2_shadowed.width) // 2
    pos2_y = card_y + 100 + (card_h - 100 - im2_shadowed.height) // 2
    canvas.paste(im2_shadowed, (pos2_x, pos2_y), im2_shadowed)

    # Connectivity Arrow
    arr_y = card_y + card_h // 2
    draw_arrow(draw, (c1_x+c1_w+10, arr_y), (c2_x-10, arr_y), fill=(148, 163, 184, 200), width=10, arrow_size=30)
    
    # Badge at the bottom
    # Adding a unified footer to lock in the layout
    footer_y = card_y + card_h + 30
    draw.text((W//2, footer_y), "The extracted cellular features flow sequentially into the Multi-Modal Integration Engine to power 92% accurate predictions.", 
              font=get_font(22, True), fill=(148, 163, 184), anchor="mm")

    # Save
    canvas.convert("RGB").save(OUT_PATH, quality=98)
    print(f"Saved PNG to {OUT_PATH}")

    # Generate PPTX
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(OUT_PATH, 0, 0, width=Inches(16))
    prs.save(PPT_PATH)
    print(f"Saved PPTX to {PPT_PATH}")

if __name__ == "__main__":
    main()
