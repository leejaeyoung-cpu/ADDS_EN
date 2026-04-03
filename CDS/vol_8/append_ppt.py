import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def append_slides_to_base(base_pptx_path, output_pptx_path):
    prs = Presentation(base_pptx_path)
    
    # Custom color palette
    primary_blue = RGBColor(0, 51, 102)
    accent_teal = RGBColor(0, 153, 153)
    
    # Since the base file might not have standard layouts in the same index,
    # we can use the default blank or title and content layout from a new presentation.
    # To be safe, let's use layout 1 (Title and Content) or layout 6 (Blank).
    # layout 1 is usually Title and Content.
    try:
        layout = prs.slide_layouts[1]
    except Exception:
        layout = prs.slide_layouts[0] # Fallback
        
    # We will add a transition/divider slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "[추가 제안 슬라이드]\n디지털 의료기기(SaMD) 상용화 및 통합 파이프라인"
    title.text_frame.paragraphs[0].font.color.rgb = primary_blue
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "AI 의사결정지원 시스템 (ADDS. V1) 고도화 로드맵"

    # Slide 1: Solution ADDS. V1
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "1. 우리의 해결책: ADDS. V1"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "다중모달 기반 AI 의사결정지원 독립형 소프트웨어"
        tf.add_paragraph().text = " - 임상 + 병리(WSI) + CT + 분자진단 데이터 통합 분석"
        tf.add_paragraph().text = " - 치료 개시 전 환자별 반응성 및 조기 진행 위험 점수 도출"
        tf.add_paragraph().text = " - PDO 기반 Ex Vivo 모델을 수반한 예측 근거와 신뢰성 확보"
        
    # Slide 2: Biomarker
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "2. 차별화 포인트: PrPᶜ 신규 바이오마커 결합"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "종양 침윤 및 전이 예측 고도화"
        tf.add_paragraph().text = " - PrPᶜ 고발현 시 미세혈관 침윤 및 병기 악화와의 강력한 연관성 확인"
        tf.add_paragraph().text = " - 단순 블랙박스 AI가 아닌 기전적/생물학적(Mechanistic) 타당성 동반"

    # Slide 3: Roadmap
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "3. 디지털 의료기기(SaMD) 상용화 로드맵"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "연구용 프로토타입을 넘어선 상용화 지향 마일스톤"
        tf.add_paragraph().text = " [초기 사양 확립] 대상(KRAS 대장암), 단일 시점(치료전) 고정"
        tf.add_paragraph().text = " [인허가 준비] 품목 분류 및 등급 판정, 의료기기 제조/품질 관리체계(QS) 도입"
        tf.add_paragraph().text = " [시판 후 관리] 변경관리, 문서/형상 통제 및 CAPA/Post-Market 감시 체계 수립"

    # Slide 4: Tech Doc & Validation
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "4. 소프트웨어 기술문서 및 임상근거 패키징"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "철저한 V&V(검증/밸리데이션) 및 임상 실증 강화"
        tf.add_paragraph().text = " - 기술문서: 요구사항 명세 -> 기능/알고리즘 정의 -> DB 정의 -> 위험관리(Risk) 문서화"
        tf.add_paragraph().text = " - 임상근거 (3 Step): (1) 후향 코호트 검증 -> (2) PDO 기반 기능 검증 -> (3) 외부/전향 환경 검증"
        tf.add_paragraph().text = " - 품질 보증을 통한 인허가용 '시험성적서' 및 '임상시험보고서' 도출 준비"

    # Slide 5: Usability & Security
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "5. 사용적합성(Usability) 및 사이버보안"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "의료기기 기준 규격을 충족하는 안전한 임상 적용 환경 구축"
        tf.add_paragraph().text = " - [사용적합성] 임상의의 편의를 위한 워크플로우 통일, 입력 오류 원천 차단 설계, XAI 제공"
        tf.add_paragraph().text = " - [보안 체계] 실시간 로그 추적, 환자 데이터 종단 간 암호화, 계정/권한 관리"
        tf.add_paragraph().text = " - 안전한 정기 업데이트 관리(버전 통제) 기반 확립"

    # Slide 6: Conclusion
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = "6. 최종 목표 및 향후 비전"
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.text = "AI 정밀 종양학 시장의 선도적 레퍼런스 구축"
        tf.add_paragraph().text = " - 규제 프레임(식약처 가이드라인)을 충족하는 '독립형 의료기기 소프트웨어' 완성"
        tf.add_paragraph().text = " - 실제 임상 진료 현장 통합(MDT 지원)을 통한 불필요한 항암 최소화"
        tf.add_paragraph().text = " => 궁극적으로 맞춤형 정밀 항암 AI 솔루션의 글로벌 확장으로 도약"

    prs.save(output_pptx_path)
    print(f"Appended additional slides to {base_pptx_path} and saved as {output_pptx_path}")

if __name__ == '__main__':
    base_file = "f:\\ADDS_CDS_PART2\\임상수요형중개연구 발표평가 준비.pptx"
    out_file = "f:\\ADDS_CDS_PART2\\ADDS_V1_Pitch_Final.pptx"
    append_slides_to_base(base_file, out_file)
