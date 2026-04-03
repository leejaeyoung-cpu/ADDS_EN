import os
from PIL import Image, ImageDraw, ImageFont, ImageFilter
import math

W, H = 1920, 1080
OUT_PATH = r"f:\ADDS\CDS\ADDS_System_Infographic.png"
PPT_PATH = r"f:\ADDS\CDS\ADDS_System_Infographic.pptx"

def get_font(size, bold=False):
    font_name = "malgunbd.ttf" if bold else "malgun.ttf"
    font_path = os.path.join(os.environ.get("WINDIR", "C:\\Windows"), "Fonts", font_name)
    try:
        return ImageFont.truetype(font_path, size)
    except:
        return ImageFont.load_default()

def create_gradient_bg(w, h, c1, c2):
    base = Image.new("RGBA", (w, h), c1)
    draw = ImageDraw.Draw(base)
    for y in range(h):
        r = int(c1[0] + (c2[0] - c1[0]) * y / h)
        g = int(c1[1] + (c2[1] - c1[1]) * y / h)
        b = int(c1[2] + (c2[2] - c1[2]) * y / h)
        draw.line([(0, y), (w, y)], fill=(r, g, b, 255))
    return base

def draw_shadowed_box(canvas, bounds, fill, radius=20, shadow_blur=15, shadow_offset=(5, 5), shadow_color=(0,0,0,50)):
    shadow = Image.new("RGBA", canvas.size, (0,0,0,0))
    s_draw = ImageDraw.Draw(shadow)
    sx = [bounds[0]+shadow_offset[0], bounds[1]+shadow_offset[1],
          bounds[2]+shadow_offset[0], bounds[3]+shadow_offset[1]]
    s_draw.rounded_rectangle(sx, radius=radius, fill=shadow_color)
    shadow = shadow.filter(ImageFilter.GaussianBlur(shadow_blur))
    canvas.paste(shadow, (0,0), shadow)
    
    c_draw = ImageDraw.Draw(canvas)
    c_draw.rounded_rectangle(bounds, radius=radius, fill=fill)
    return c_draw

def draw_rounded_rect(draw, bounds, fill, radius=20, outline=None, width=0):
    draw.rounded_rectangle(bounds, radius=radius, fill=fill, outline=outline, width=width)

def draw_arrow(draw, start, end, fill, width=6, arrow_size=20):
    draw.line([start, end], fill=fill, width=width)
    angle = math.atan2(end[1]-start[1], end[0]-start[0])
    p1 = (end[0] - arrow_size * math.cos(angle - math.pi/6),
          end[1] - arrow_size * math.sin(angle - math.pi/6))
    p2 = (end[0] - arrow_size * math.cos(angle + math.pi/6),
          end[1] - arrow_size * math.sin(angle + math.pi/6))
    draw.polygon([end, p1, p2], fill=fill)

def main():
    print("Generating ADDS Infographic...")
    # Soft clinical background
    canvas = create_gradient_bg(W, H, (244, 248, 252), (230, 238, 248))
    draw = ImageDraw.Draw(canvas)

    # Header
    draw_rounded_rect(draw, [0, 0, W, 130], fill=(26, 54, 93, 255), radius=0)
    draw.text((60, 30), "ADDS: AI-Driven Decision Support System", font=get_font(52, True), fill=(255, 255, 255))
    draw.text((W - 750, 50), "End-to-End Precision Oncology Architecture", font=get_font(30, False), fill=(160, 186, 224))

    # Core Layout Coordinates
    box_y = 180
    box_h = 820
    
    f_title = get_font(34, True)
    f_sub = get_font(24, True)
    f_body = get_font(20, False)

    # ─────────────────────────────────────────────────────────────────
    # TIER 1: Source Data (LEFT)
    # ─────────────────────────────────────────────────────────────────
    b1_x = 60
    b1_w = 460
    draw_shadowed_box(canvas, [b1_x, box_y, b1_x+b1_w, box_y+box_h], fill=(255, 255, 255, 255))
    
    # Tier 1 Header
    draw_rounded_rect(draw, [b1_x, box_y, b1_x+b1_w, box_y+90], fill=(49, 130, 206, 255), radius=20)
    draw_rounded_rect(draw, [b1_x, box_y+70, b1_x+b1_w, box_y+90], fill=(49, 130, 206, 255), radius=0)
    draw.text((b1_x+40, box_y+25), "Tier 1: Source Streams", font=f_title, fill=(255,255,255))

    inputs = [
        ("Genomic Database", "KRAS, MSI, TP53\nMutational Burden (TMB)", (235, 244, 255), (49, 130, 206)),
        ("Cellular Analysis (WSI)", "H&E Imaging\nMorphology, Ki-67", (230, 255, 250), (49, 151, 149)),
        ("Imaging Data (CT)", "3D Tumor Volume\nVascularity & TNM", (255, 245, 245), (229, 62, 62)),
        ("Clinical EMR", "Age, ECOG PS\nPrior Treatment History", (255, 250, 240), (221, 107, 32))
    ]
    
    cy = box_y + 130
    for title, desc, bg, fg in inputs:
        draw_rounded_rect(draw, [b1_x+30, cy, b1_x+b1_w-30, cy+140], fill=bg, outline=fg, width=2)
        draw.text((b1_x+50, cy+25), title, font=get_font(26, True), fill=fg)
        draw.text((b1_x+50, cy+70), desc, font=get_font(22, False), fill=(45, 55, 72))
        cy += 165

    # ─────────────────────────────────────────────────────────────────
    # TIER 2: Integration Engine (CENTER)
    # ─────────────────────────────────────────────────────────────────
    b2_x = b1_x + b1_w + 60
    b2_w = 760
    
    # Glowing back box
    draw_shadowed_box(canvas, [b2_x, box_y, b2_x+b2_w, box_y+box_h], fill=(26, 32, 44, 255), shadow_blur=25, shadow_color=(49, 130, 206, 120))
    
    # Tier 2 Header
    draw_rounded_rect(draw, [b2_x, box_y, b2_x+b2_w, box_y+90], fill=(43, 108, 176, 255), radius=20)
    draw_rounded_rect(draw, [b2_x, box_y+70, b2_x+b2_w, box_y+90], fill=(43, 108, 176, 255), radius=0)
    draw.text((b2_x+40, box_y+25), "Tier 2: AI Integration Engine", font=f_title, fill=(255,255,255))

    # Part A: Cellpose Pipeline
    cy = box_y + 130
    draw_rounded_rect(draw, [b2_x+30, cy, b2_x+b2_w-30, cy+300], fill=(45, 55, 72, 255), outline=(99, 179, 237, 255), width=2)
    draw.text((b2_x+50, cy+20), "A. Cellpose Deep Learning Segmentation", font=get_font(28, True), fill=(99, 179, 237))
    draw.text((b2_x+50, cy+70), "▶ 512x512 Tiling & Color Normalization\n▶ Vector Flow Field & Cell Probability Map (U-Net)\n▶ Feature Extraction (Area, Circularity, Atypia)", font=get_font(22, False), fill=(226, 232, 240))
    
    # Metrics
    my = cy + 200
    draw_rounded_rect(draw, [b2_x+50, my, b2_x+260, my+80], fill=(26, 32, 44, 255), outline=(104, 211, 145, 255), width=2)
    draw.text((b2_x+65, my+15), "Accuracy\n95.2%", font=get_font(20, True), fill=(104, 211, 145))
    
    draw_rounded_rect(draw, [b2_x+280, my, b2_x+490, my+80], fill=(26, 32, 44, 255), outline=(246, 173, 85, 255), width=2)
    draw.text((b2_x+295, my+15), "GPU Process\n3 seconds", font=get_font(20, True), fill=(246, 173, 85))
    
    draw_rounded_rect(draw, [b2_x+510, my, b2_x+710, my+80], fill=(26, 32, 44, 255), outline=(252, 129, 129, 255), width=2)
    draw.text((b2_x+525, my+15), "Viability\n73.5%", font=get_font(20, True), fill=(252, 129, 129))

    # Part B: Therapy Selection
    cy += 330
    draw_rounded_rect(draw, [b2_x+30, cy, b2_x+b2_w-30, cy+320], fill=(45, 55, 72, 255), outline=(159, 122, 234, 255), width=2)
    draw.text((b2_x+50, cy+20), "B. Polypharmacology Selection Ensemble", font=get_font(28, True), fill=(179, 136, 255))
    draw.text((b2_x+50, cy+70), "▶ Algorithms: Random Forest + CNN Ensemble\n▶ Pipeline: Extraction → Interaction Modeling → Efficacy\n▶ Synergy Prediction cross-mapped with NCCN", font=get_font(22, False), fill=(226, 232, 240))
    
    my2 = cy + 210
    draw_rounded_rect(draw, [b2_x+50, my2, b2_x+b2_w-50, my2+90], fill=(26, 32, 44, 255), outline=(236, 201, 75, 255), width=2)
    draw.text((b2_x+80, my2+15), "Prediction Accuracy: 92%   |   Model Performance: R² = 0.97", font=get_font(22, True), fill=(236, 201, 75))
    draw.text((b2_x+80, my2+50), "★ Publication Quality Statistical Significance (p < 0.001)", font=get_font(18, False), fill=(226, 232, 240))

    # ─────────────────────────────────────────────────────────────────
    # TIER 3: Outputs (RIGHT)
    # ─────────────────────────────────────────────────────────────────
    b3_x = b2_x + b2_w + 60
    b3_w = W - b3_x - 60
    draw_shadowed_box(canvas, [b3_x, box_y, b3_x+b3_w, box_y+box_h], fill=(255, 255, 255, 255))
    
    # Tier 3 Header
    draw_rounded_rect(draw, [b3_x, box_y, b3_x+b3_w, box_y+90], fill=(56, 161, 105, 255), radius=20)
    draw_rounded_rect(draw, [b3_x, box_y+70, b3_x+b3_w, box_y+90], fill=(56, 161, 105, 255), radius=0)
    draw.text((b3_x+40, box_y+25), "Tier 3: Actionable Outputs", font=f_title, fill=(255,255,255))

    cy = box_y + 130
    
    # Rec 1
    draw_rounded_rect(draw, [b3_x+30, cy, b3_x+b3_w-30, cy+150], fill=(240, 255, 244, 255), outline=(72, 187, 120, 255), width=2)
    draw.text((b3_x+50, cy+25), "Optimal: FOLFOX + Bev", font=get_font(26, True), fill=(34, 84, 61))
    draw.text((b3_x+50, cy+80), "Efficacy: 78% | Confidence: 94%", font=get_font(24, False), fill=(39, 103, 73))
    
    cy += 180
    # Rec 2
    draw_rounded_rect(draw, [b3_x+30, cy, b3_x+b3_w-30, cy+150], fill=(255, 250, 240, 255), outline=(237, 137, 54, 255), width=2)
    draw.text((b3_x+50, cy+25), "Alt: CAPOX + Bevacizumab", font=get_font(26, True), fill=(140, 58, 0))
    draw.text((b3_x+50, cy+80), "Efficacy: 82% | Confidence: 87%", font=get_font(24, False), fill=(192, 86, 33))

    cy += 180
    # Automated Dual Reports
    draw_rounded_rect(draw, [b3_x+30, cy, b3_x+b3_w-30, cy+150], fill=(235, 244, 255, 255), outline=(49, 130, 206, 255), width=2)
    draw.text((b3_x+50, cy+30), "Dual Document Generation", font=get_font(26, True), fill=(43, 108, 176))
    draw.text((b3_x+50, cy+85), "1. Medical Doctor (MD) Report\n2. Patient Visual Dashboard", font=get_font(24, False), fill=(43, 108, 176))
    
    # ── Arrows connecting the tiers ──
    # T1 -> T2
    arrow_y = H//2 + 40
    draw_arrow(draw, (b1_x+b1_w+5, arrow_y-150), (b2_x-5, arrow_y-150), fill=(160, 174, 192), width=8, arrow_size=25)
    draw_arrow(draw, (b1_x+b1_w+5, arrow_y+200), (b2_x-5, arrow_y+200), fill=(160, 174, 192), width=8, arrow_size=25)
    
    # T2 -> T3
    draw_arrow(draw, (b2_x+b2_w+5, arrow_y), (b3_x-5, arrow_y), fill=(72, 187, 120), width=10, arrow_size=25)

    # Save outputs
    canvas.save(OUT_PATH)
    print(f"Saved PNG to {OUT_PATH}")

    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(OUT_PATH, 0, 0, width=Inches(16))
    prs.save(PPT_PATH)
    print(f"Saved PPTX to {PPT_PATH}")

if __name__ == "__main__":
    main()
