import os
from PIL import Image, ImageDraw, ImageFont, ImageFilter

W, H = 1920, 1080
OUT_PATH = r"f:\ADDS\CDS\NanoBanana2_Simplified_Pipeline.png"
PPT_PATH = r"f:\ADDS\CDS\NanoBanana2_Simplified_Pipeline.pptx"

def get_font(size, bold=False):
    font_name = "malgunbd.ttf" if bold else "malgun.ttf"
    font_path = os.path.join(os.environ.get("WINDIR", "C:\\Windows"), "Fonts", font_name)
    try:
        return ImageFont.truetype(font_path, size)
    except:
        return ImageFont.load_default()

def draw_rounded_rect(draw, bounds, fill, radius=20, outline=None, width=0):
    draw.rounded_rectangle(bounds, radius=radius, fill=fill, outline=outline, width=width)

def draw_shadowed_box(canvas, bounds, fill, radius=20, shadow_blur=15, shadow_offset=(5, 5), shadow_color=(0,0,0,60)):
    shadow = Image.new("RGBA", canvas.size, (0,0,0,0))
    s_draw = ImageDraw.Draw(shadow)
    sx = [b + offset for b, offset in zip(bounds, shadow_offset*2)]
    s_draw.rounded_rectangle(sx, radius=radius, fill=shadow_color)
    shadow = shadow.filter(ImageFilter.GaussianBlur(shadow_blur))
    canvas.paste(shadow, (0,0), shadow)
    
    c_draw = ImageDraw.Draw(canvas)
    c_draw.rounded_rectangle(bounds, radius=radius, fill=fill)
    return c_draw

def create_gradient_bg(w, h, c1, c2):
    base = Image.new("RGBA", (w, h), c1)
    draw = ImageDraw.Draw(base)
    for y in range(h):
        r = int(c1[0] + (c2[0] - c1[0]) * y / h)
        g = int(c1[1] + (c2[1] - c1[1]) * y / h)
        b = int(c1[2] + (c2[2] - c1[2]) * y / h)
        draw.line([(0, y), (w, y)], fill=(r, g, b, 255))
    return base

def draw_arrow(draw, start, end, fill, width=8, arrow_size=20):
    draw.line([start, end], fill=fill, width=width)
    import math
    angle = math.atan2(end[1]-start[1], end[0]-start[0])
    p1 = (end[0] - arrow_size * math.cos(angle - math.pi/6),
          end[1] - arrow_size * math.sin(angle - math.pi/6))
    p2 = (end[0] - arrow_size * math.cos(angle + math.pi/6),
          end[1] - arrow_size * math.sin(angle + math.pi/6))
    draw.polygon([end, p1, p2], fill=fill)

def main():
    print("Generating Simplified Pipeline Infographic...")
    # Background: clean modern light gradient
    canvas = create_gradient_bg(W, H, (248, 250, 252), (226, 232, 240))
    draw = ImageDraw.Draw(canvas)

    # Header
    draw_rounded_rect(draw, [0, 0, W, 140], fill=(15, 23, 42, 255), radius=0)
    draw.text((80, 40), "ADDS: Unified Precision Oncology AI Pipeline", font=get_font(52, True), fill=(255, 255, 255))
    draw.text((W - 600, 60), "End-to-End Seamless AI Workflow", font=get_font(28, False), fill=(148, 163, 184))

    # Fonts
    f_title = get_font(38, True)
    f_sub = get_font(24, True)
    f_body = get_font(22, False)
    f_metric = get_font(32, True)

    box_y = 220
    box_h = 750

    # ─────────────────────────────────────────────────────────────────
    # Box 1: Source Data Streams (Left)
    # ─────────────────────────────────────────────────────────────────
    b1_x = 80
    b1_w = 480
    draw_shadowed_box(canvas, [b1_x, box_y, b1_x+b1_w, box_y+box_h], fill=(255, 255, 255, 255))
    
    # Header 1
    draw_rounded_rect(draw, [b1_x, box_y, b1_x+b1_w, box_y+100], fill=(56, 189, 248, 255), radius=20)
    draw_rounded_rect(draw, [b1_x, box_y+80, b1_x+b1_w, box_y+100], fill=(56, 189, 248, 255), radius=0) # flatten bottom
    draw.text((b1_x + 30, box_y + 30), "1. Source Data Streams", font=f_title, fill=(255, 255, 255))

    cy = box_y + 140
    data_items = [
        ("🧬 Genomic Data", "KRAS, TP53, MSI, TMB Score\nSequence DB matching"),
        ("🔬 Cellular Analysis (WSI)", "Cell morphology, Ki-67 index\nH&E Patch features"),
        ("🩻 Imaging Data (CT)", "3D Tumor volume calculation\nVascularity & TNM Stage"),
        ("📋 Clinical EMR Data", "Patient history, ECOG PS\nPrior lines of treatment")
    ]
    for title, desc in data_items:
        draw_rounded_rect(draw, [b1_x+30, cy, b1_x+b1_w-30, cy+125], fill=(241, 245, 249, 255), outline=(203, 213, 225, 255), width=2)
        draw.text((b1_x + 50, cy + 20), title, font=f_sub, fill=(15, 23, 42))
        draw.text((b1_x + 50, cy + 60), desc, font=f_body, fill=(71, 85, 105))
        cy += 145

    # Arrow 1->2
    draw_arrow(draw, (b1_x+b1_w+10, box_y + box_h//2), (b1_x+b1_w+70, box_y + box_h//2), fill=(148, 163, 184), width=12, arrow_size=30)

    # ─────────────────────────────────────────────────────────────────
    # Box 2: ADDS Core (Center)
    # ─────────────────────────────────────────────────────────────────
    b2_x = b1_x + b1_w + 80
    b2_w = 640
    # Make it glowing/darker for emphasis
    draw_shadowed_box(canvas, [b2_x, box_y, b2_x+b2_w, box_y+box_h], fill=(15, 23, 42, 255), shadow_blur=30, shadow_color=(37, 99, 235, 100))
    
    draw_rounded_rect(draw, [b2_x, box_y, b2_x+b2_w, box_y+100], fill=(37, 99, 235, 255), radius=20)
    draw_rounded_rect(draw, [b2_x, box_y+80, b2_x+b2_w, box_y+100], fill=(37, 99, 235, 255), radius=0)
    draw.text((b2_x + 30, box_y + 30), "2. ADDS Core AI Engine", font=f_title, fill=(255, 255, 255))

    cy = box_y + 130
    
    # AI Engine Part 1
    draw.text((b2_x + 40, cy), "A. Multimodal Cellular Segmentation", font=f_sub, fill=(96, 165, 250))
    cy += 40
    draw.text((b2_x + 40, cy), "▶ Vector Flow Fields & Cellpose Deep Learning", font=f_body, fill=(203, 213, 225))
    cy += 30
    draw.text((b2_x + 40, cy), "▶ High-fidelity extraction of Atypia Index & N/C Ratio", font=f_body, fill=(203, 213, 225))
    
    cy += 60
    # Metrics horizontally
    draw_rounded_rect(draw, [b2_x+30, cy, b2_x+b2_w//2-10, cy+90], fill=(30, 41, 59, 255), outline=(56, 189, 248, 255), width=2)
    draw.text((b2_x+50, cy+15), "Segmentation Acc.", font=get_font(18), fill=(148, 163, 184))
    draw.text((b2_x+50, cy+40), "95.2%", font=f_metric, fill=(56, 189, 248))

    draw_rounded_rect(draw, [b2_x+b2_w//2+10, cy, b2_x+b2_w-30, cy+90], fill=(30, 41, 59, 255), outline=(16, 185, 129, 255), width=2)
    draw.text((b2_x+b2_w//2+30, cy+15), "GPU Processing", font=get_font(18), fill=(148, 163, 184))
    draw.text((b2_x+b2_w//2+30, cy+40), "3.0 sec (15x UP)", font=f_metric, fill=(16, 185, 129))

    cy += 140
    # AI Engine Part 2
    draw.text((b2_x + 40, cy), "B. Polypharmacology Synergy Prediction", font=f_sub, fill=(96, 165, 250))
    cy += 40
    draw.text((b2_x + 40, cy), "▶ LLM + Random Forest Hybrid Architecture", font=f_body, fill=(203, 213, 225))
    cy += 30
    draw.text((b2_x + 40, cy), "▶ Cross-matching with NCCN & Target Interactions", font=f_body, fill=(203, 213, 225))
    
    cy += 60
    draw_rounded_rect(draw, [b2_x+30, cy, b2_x+b2_w-30, cy+100], fill=(30, 41, 59, 255), outline=(245, 158, 11, 255), width=2)
    draw.text((b2_x+50, cy+20), "Model Performance", font=get_font(20), fill=(148, 163, 184))
    draw.text((b2_x+50, cy+50), "R² = 0.97 (p < 0.001)", font=get_font(28, True), fill=(245, 158, 11))
    draw.text((b2_x+360, cy+45), "★ Publication Quality", font=get_font(20, True), fill=(255, 255, 255))
    
    # Arrow 2->3
    draw_arrow(draw, (b2_x+b2_w+10, box_y + box_h//2), (b2_x+b2_w+70, box_y + box_h//2), fill=(148, 163, 184), width=12, arrow_size=30)

    # ─────────────────────────────────────────────────────────────────
    # Box 3: Output (Right)
    # ─────────────────────────────────────────────────────────────────
    b3_x = b2_x + b2_w + 80
    b3_w = W - b3_x - 80
    draw_shadowed_box(canvas, [b3_x, box_y, b3_x+b3_w, box_y+box_h], fill=(255, 255, 255, 255))
    
    draw_rounded_rect(draw, [b3_x, box_y, b3_x+b3_w, box_y+100], fill=(16, 185, 129, 255), radius=20)
    draw_rounded_rect(draw, [b3_x, box_y+80, b3_x+b3_w, box_y+100], fill=(16, 185, 129, 255), radius=0)
    draw.text((b3_x + 30, box_y + 30), "3. Cocktail Recommendations", font=f_title, fill=(255, 255, 255))

    cy = box_y + 130
    
    # Recommendation Box
    draw_rounded_rect(draw, [b3_x+20, cy, b3_x+b3_w-20, cy+200], fill=(236, 253, 245, 255), outline=(16, 185, 129, 255), width=3)
    draw_rounded_rect(draw, [b3_x+b3_w-150, cy, b3_x+b3_w-20, cy+40], fill=(16, 185, 129, 255), radius=10)
    draw.text((b3_x+b3_w-125, cy+7), "1st Choice", font=get_font(18, True), fill=(255,255,255))
    
    draw.text((b3_x + 40, cy + 20), "FOLFOX\n+ Bevacizumab", font=get_font(26, True), fill=(6, 78, 59))
    draw.text((b3_x + 40, cy + 90), "• Predicted Efficacy: 89.5%\n• Conf. Score: 94%\n• Side Effect Risk: Moderate", font=f_body, fill=(4, 120, 87))

    cy += 230
    # Alternative Box
    draw_rounded_rect(draw, [b3_x+20, cy, b3_x+b3_w-20, cy+180], fill=(255, 251, 235, 255), outline=(245, 158, 11, 255), width=2)
    draw.text((b3_x + 40, cy + 20), "Alternative:\nFOLFIRI + Bev", font=get_font(24, True), fill=(180, 83, 9))
    draw.text((b3_x + 40, cy + 90), "• Pred. Efficacy: 85.2%\n• Risk: Moderate-High", font=f_body, fill=(146, 64, 14))

    cy += 210
    # Dual Output
    draw_rounded_rect(draw, [b3_x+20, cy, b3_x+b3_w-20, cy+130], fill=(241, 245, 249, 255), outline=(148, 163, 184, 255), width=2)
    draw.text((b3_x + 40, cy + 15), "Dual Output Reports", font=get_font(22, True), fill=(15, 23, 42))
    draw.text((b3_x + 40, cy + 50), "1) Patient Dashboard\n   (Visual & Intuitive)\n2) Clinical MD Report", font=get_font(20, False), fill=(71, 85, 105))


    # Save
    canvas.save(OUT_PATH)
    print(f"Saved PNG to {OUT_PATH}")

    # Generate PPTX
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(OUT_PATH, 0, 0, width=Inches(16))
    prs.save(PPT_PATH)
    print(f"Saved PPTX to {PPT_PATH}")

if __name__ == "__main__":
    main()
