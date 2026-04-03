import os
import math
from PIL import Image, ImageDraw, ImageFont, ImageFilter

# ─── Settings ──────────────
WIDTH, HEIGHT = 1920, 1080
BG_COLOR_1 = (13, 19, 33)        # #0D1321
BG_COLOR_2 = (29, 45, 68)        # #1D2D44
ACCENT_BLUE = (37, 99, 235)      # #2563eb
ACCENT_GREEN = (16, 185, 129)    # #10b981
TEXT_WHITE = (248, 250, 252)
TEXT_GRAY = (148, 163, 184)

# Output path
OUT_PATH = r"f:\ADDS\CDS\ADDS_Dual_Report_Showcase.png"

# Inputs
IMG_PATIENT = r"f:\ADDS\CDS\NanoBanana2_Patient_Dashboard_FINAL_Page_1.png"
IMG_DOCTOR = r"f:\ADDS\CDS\NanoBanana2_Doctor_Report_FINAL_Page_1.png"

def get_font(size, bold=False):
    font_name = "malgunbd.ttf" if bold else "malgun.ttf"
    font_path = os.path.join(os.environ.get("WINDIR", "C:\\Windows"), "Fonts", font_name)
    try:
        return ImageFont.truetype(font_path, size)
    except:
        return ImageFont.load_default()

def create_gradient_bg(w, h, c1, c2):
    base = Image.new("RGBA", (w, h), c1)
    draw = ImageDraw.Draw(base)
    for y in range(h):
        r = int(c1[0] + (c2[0] - c1[0]) * y / h)
        g = int(c1[1] + (c2[1] - c1[1]) * y / h)
        b = int(c1[2] + (c2[2] - c1[2]) * y / h)
        draw.line([(0, y), (w, y)], fill=(r, g, b, 255))
    return base

def add_drop_shadow(image, offset=(15, 15), blur_radius=20, shadow_color=(0, 0, 0, 150)):
    shadow = Image.new('RGBA', image.size, shadow_color)
    # create a padded base for the shadow
    padded_size = (image.width + offset[0] + blur_radius*2,
                   image.height + offset[1] + blur_radius*2)
    base = Image.new('RGBA', padded_size, (0,0,0,0))
    base.paste(shadow, (blur_radius + offset[0], blur_radius + offset[1]))
    base = base.filter(ImageFilter.GaussianBlur(blur_radius))
    # paste real image over it
    base.paste(image, (blur_radius, blur_radius), image)
    return base

def draw_abstract_shapes(draw, w, h):
    # draw some glowing circles in the background
    for x, y, r, alpha in [
        (200, 200, 400, 30),
        (1700, 800, 500, 20),
        (960, 500, 600, 10),
    ]:
        bbox = [x-r, y-r, x+r, y+r]
        draw.ellipse(bbox, fill=(37, 99, 235, alpha))

def draw_network_nodes(draw):
    # Draws a small connected network representing AI analysis
    nodes = [(150, 850), (250, 750), (350, 880), (450, 700), (450, 850)]
    for i in range(len(nodes)):
        for j in range(i+1, len(nodes)):
            if math.hypot(nodes[i][0]-nodes[j][0], nodes[i][1]-nodes[j][1]) < 200:
                draw.line([nodes[i], nodes[j]], fill=(37, 99, 235, 100), width=2)
    for nx, ny in nodes:
        draw.ellipse([nx-8, ny-8, nx+8, ny+8], fill=(16, 185, 129, 200))


def main():
    print("Generating Showcase...")
    canvas = create_gradient_bg(WIDTH, HEIGHT, BG_COLOR_1, BG_COLOR_2)
    draw = ImageDraw.Draw(canvas, "RGBA")
    
    draw_abstract_shapes(draw, WIDTH, HEIGHT)
    draw_network_nodes(draw)
    
    # ── Title Text ──
    f_title = get_font(52, True)
    f_sub = get_font(28, False)
    
    draw.text((80, 80), "NanoBanana 2.0: Dual Output Reporting System", font=f_title, fill=TEXT_WHITE)
    draw.text((80, 150), "Single Engine, Tailored Insights: Translating AI precision into patient-friendly dashboards & clinical expert reports.", font=f_sub, fill=TEXT_GRAY)

    # ── Load and Resize Document Images ──
    try:
        img_pat = Image.open(IMG_PATIENT).convert("RGBA")
        img_doc = Image.open(IMG_DOCTOR).convert("RGBA")
    except Exception as e:
        print(f"Error loading images: {e}")
        return

    # Resize Patient Dashboard (Landscape A4)
    target_w_pat = 1000
    ratio_pat = target_w_pat / float(img_pat.width)
    img_pat_res = img_pat.resize((target_w_pat, int(img_pat.height * ratio_pat)), Image.LANCZOS)
    
    # Resize Doctor Report (Portrait A4)
    target_w_doc = 650
    ratio_doc = target_w_doc / float(img_doc.width)
    img_doc_res = img_doc.resize((target_w_doc, int(img_doc.height * ratio_doc)), Image.LANCZOS)
    
    # Add beautiful drop shadows
    pat_shadowed = add_drop_shadow(img_pat_res, offset=(10, 20), blur_radius=25, shadow_color=(0, 0, 0, 180))
    doc_shadowed = add_drop_shadow(img_doc_res, offset=(10, 20), blur_radius=25, shadow_color=(0, 0, 0, 180))
    
    # ── Paste onto canvas ──
    # Left: Patient Dashboard
    pat_x = 80
    pat_y = 300
    canvas.paste(pat_shadowed, (pat_x, pat_y), pat_shadowed)
    
    # Right: Doctor Report (slightly overlapping or layered)
    doc_x = 1150
    doc_y = 120
    canvas.paste(doc_shadowed, (doc_x, doc_y), doc_shadowed)

    # ── Annotations ──
    f_anno_title = get_font(36, True)
    f_anno_text = get_font(22, False)

    # Annotation for Patient
    draw.text((pat_x + 40, pat_y - 60), "1. 환자용 안심 대시보드 (Patient-Friendly)", font=f_anno_title, fill=(16, 185, 129))
    draw.text((pat_x + 40, pat_y + pat_shadowed.height + 10), "▶ 한눈에 들어오는 3-패널 인포그래픽", font=f_anno_text, fill=TEXT_WHITE)
    draw.text((pat_x + 40, pat_y + pat_shadowed.height + 40), "▶ 쉬운 용어, 부작용 모니터링, AI 추천 근거의 시각화", font=f_anno_text, fill=TEXT_GRAY)

    # Annotation for Doctor
    draw.text((doc_x + 60, doc_y - 60), "2. 의사용 상세 리포트 (Clinical Expert)", font=f_anno_title, fill=(37, 99, 235))
    draw.text((doc_x + 60, doc_y + doc_shadowed.height + 10), "▶ 전통적 임상 리포트 포맷 준수", font=f_anno_text, fill=TEXT_WHITE)
    draw.text((doc_x + 60, doc_y + doc_shadowed.height + 40), "▶ 17D 텐서 분석 및 정밀 종양 부피, 감수성(%) 수치 제공", font=f_anno_text, fill=TEXT_GRAY)

    # Center floating badge / connecting line
    center_cx = 1150
    center_cy = int(pat_y + pat_shadowed.height / 2)
    # Draw a line connecting the two domains
    draw.line([(pat_x + pat_shadowed.width - 40, center_cy), (doc_x + 50, center_cy)], fill=(255, 255, 255, 60), width=4)
    
    # Badge
    badge_w, badge_h = 240, 60
    badge_x = center_cx - badge_w//2 - 40
    badge_y = center_cy - badge_h//2
    draw.rounded_rectangle([badge_x, badge_y, badge_x + badge_w, badge_y + badge_h], radius=30, fill=(37, 99, 235, 220), outline=TEXT_WHITE, width=2)
    draw.text((badge_x + 25, badge_y + 12), "One AI, Dual Output", font=get_font(20, True), fill=TEXT_WHITE)

    # Save PNG
    canvas.save(OUT_PATH)
    print(f"Awesome showcase saved to: {OUT_PATH}")

    # Generate PPTX
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(16), Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(OUT_PATH, 0, 0, width=Inches(16))
    ppt_path = r"f:\ADDS\CDS\ADDS_Dual_Report_Showcase.pptx"
    prs.save(ppt_path)
    print(f"Presentation saved to: {ppt_path}")

if __name__ == "__main__":
    main()
