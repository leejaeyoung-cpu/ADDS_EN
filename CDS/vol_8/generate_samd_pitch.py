import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import datetime

def rgb(r, g, b):
    return RGBColor(r, g, b)

# Colors
COLOR_DARK_BLUE = rgb(18, 35, 66)
COLOR_TEAL = rgb(25, 160, 160)
COLOR_MINT = rgb(72, 209, 204)
COLOR_WHITE = rgb(255, 255, 255)
COLOR_GRAY = rgb(100, 100, 100)

def create_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BLUE

    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle = slide.placeholders[1]
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_MINT
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    return slide

def create_content_slide(prs, title_text, bullet_points, emphasis_box=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK_BLUE
    title.text_frame.paragraphs[0].font.bold = True
    
    # Content body
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for bp in bullet_points:
        p = tf.add_paragraph()
        p.text = bp
        p.font.size = Pt(20)
        p.font.color.rgb = rgb(50, 50, 50)
        p.level = 0
        p.space_after = Pt(14)
        
    # Add emphasis box if provided
    if emphasis_box:
        left = Inches(1)
        top = Inches(5.5)
        width = Inches(8)
        height = Inches(1.2)
        
        shape = slide.shapes.add_shape(
            1, left, top, width, height # 1 is MSO_SHAPE.RECTANGLE
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(240, 248, 255) # Light Alice Blue
        shape.line.color.rgb = COLOR_TEAL
        shape.line.width = Pt(2)
        
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = emphasis_box
        p.font.color.rgb = COLOR_DARK_BLUE
        p.font.bold = True
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
        
    return slide

def build_presentation(out_path):
    prs = Presentation()
    
    # Slide 1: Title
    create_title_slide(prs, 
        "KRAS 변이 대장암 기반 AI 의사결정지원 의료기술 (ADDS. V1) 실용화", 
        "단순 분석 모델을 넘어 '독립형 AI 디지털의료기기(SaMD)' 제품화 및 확산 전략\n\n임상수요형 중개연구 발표평가\n" + datetime.datetime.now().strftime("%Y-%m-%d"))

    # Slide 2: Background
    create_content_slide(prs, 
        "01. 임상 미충족 수요 및 ADDS. V1의 연구 성과 요약", 
        [
            "임상 현장 한계: 대장암 환자별 종양 이질성으로 표준요법 적용 시 치료 반응 편차 큼",
            "ADDS V1 연구 성과: 임상·WSI·CT 융합 분석으로 치료 반응/조기 진행을 예측하는 '항암 의사결정 지원 시스템' 기반 확립 완료",
            "연구에서 실질로: 예측 모델 개발 완료에 머무르지 않고 실제 의료 현장에서 즉시 적용 가능한 '신뢰 기반 의료기기'로의 고도화가 필수적"
        ],
        "핵심 가치: 연구 환경의 우수한 '정확도'를 실제 임상의 '신뢰도'로 치환하기 위한 체계적인 상용화 전략 필요")

    # Slide 3: Vision
    create_content_slide(prs, 
        "02. 연구용 모델에서 독립형 AI 디지털의료기기(SaMD)로의 도약", 
        [
            "기존: 개별 연구 단위의 산점된 딥러닝 결과 도출 (연구 성과 중심)",
            "목표: 규제(Regulatory)가 반영된 독립형 소프트웨어 의료기기(SaMD) 전환",
            "접근 전략: 식약처 인허가 가이드라인에 맞춘 단계적 포트폴리오(Productization)",
            "파급 효과: 단일 솔루션이 아닌 의료진의 표준화된 치료 도구로 시장 안착 지향"
        ],
        "[Vision] 병원 시스템(EMR/CDSS)과 완벽히 연동되는 'AI-assisted 정밀 항암 솔루션' 구축")

    # Slide 4: Regulatory Strategy
    create_content_slide(prs, 
        "03. SaMD 해당성 검토 및 품목/등급 분류 설계", 
        [
            "의료기기 해당성: 환자의 진료 기록(CT, 병리, 임상 데이터)을 정량 분석하여 진단 및 치료 계획(항암 병용) 정보를 제공하므로 SaMD 해당",
            "품목 분류: 「의료영상진단보조소프트웨어」 또는 「방사선암치료계획소프트웨어」 기준 검토",
            "위험 등급: 환자 맞춤형 치료 전략을 직접적으로 제안하므로 3등급(잠재적 위해성 높음) 판정 예상 ➔ 엄격한 임상 검증 요구됨",
            "인허가/인증 경로: 혁신의료기기 지정 신청 및 통합심사(신의료기술평가) 패스트트랙 활용 조기 시장진입 모색"
        ],
        "Regulatory Action: 식약처 등급 분류 선행 심사 및 혁신의료기기 통합심사 트랙 준비")

    # Slide 5: Quality & Manufacturing
    create_content_slide(prs, 
        "04. 의료기기 제조 및 품질관리체계(GMP/QS) 구축", 
        [
            "디지털 의료기기 특화 가이드라인: 국내 GMP(KGMP) 및 ISO 13485 규격 기반 체계 구축 설계",
            "제조 주체: 소프트웨어 단독 제품이므로 개발 기관(본 연구팀 연계) 주체의 멸균/위생이 제외된 SW 설계 기반 QMS 문서 정립",
            "문서 통제 및 파생 관리: 형상관리(Configuration Management), 버전 변경, 코드 배포 과정 통제 SOP 확립",
            "설계 이력 및 CAPA: 개발 단계별 설계 이력 파일(DHF) 유지 및 시판 전/후 이상 조치(CAPA) 매뉴얼화"
        ],
        "제조 QS 방향: 순수 소프트웨어에 특화된 Agile 품질 통제 및 무결성 기반 지속적 업데이트 체계 확립")

    # Slide 6: Tech Docs
    create_content_slide(prs, 
        "05. 소프트웨어 기술문서 패키지화", 
        [
            "요구사항 명세성: 의료진의 임상적 수요(조기 진행 예측 등)를 SW 기능 요구사항으로 치환 문서화",
            "알고리즘 명세서: Cellpose 기반 3D 분할, 하이브리드 CT 엔진, 딥러닝(Swin-UNETR) 엔진 등 로직 설명서",
            "위험관리: ISO 14971 기반의 오분류(위양성/위음성), SW 중단, 데이터 누락에 대한 식별 및 통제 방안 파일화",
            "V&V (검증 및 타당성 확인): 시스템 모듈 단위 시험 성적서(제 3자 공인 혹은 내부 엄격 체계) 발급"
        ],
        "Tech Package: 식약처 가이드라인 기반의 AI 알고리즘 설명서, 데이터셋 정의서 및 위험관리 문서화 완비")

    # Slide 7: Clinical Evidence
    create_content_slide(prs, 
        "06. 임상근거 (Clinical Evidence) 확립 단계적 전략", 
        [
            "Phase 1. 후향적 근거: 기 확보된 내부 및 외부 코호트, 멀티센터 데이터를 통한 AUC/Dice Score 통계적 유의성 확보",
            "Phase 2. 체외 효능: 환자 유래 오가노이드(PDO)를 활용한 예측 결과-실제 약물 반응(IC50) 쌍의 실험적 확인",
            "Phase 3. 전향적 관찰: 확증적 임상시험을 대체/보완할 전향적 실사용 환경 파일럿 적용 (초진 시 적용 가치 관찰)",
            "목표치: 분할 Dice > 0.8, 반응 예측 AUC > 0.85 임상적 임계 수치 만족 증명서 발급"
        ],
        "실증(Evidence): 후향 코호트 + PDO 검증 + 다기관 전향적 관찰의 3-Tier 임상 안전성 검증 파이프라인")

    # Slide 8: Usability & Security
    create_content_slide(prs, 
        "07. 사용적합성(Usability) 및 사이버보안 체계 확립", 
        [
            "사용적합성(IEC 62366): 의료인의 결과 판독 시 휴먼 에러를 방지할 UI 설계, 색상, 배치 및 혼동 요소 제거 시험",
            "오류 방지 설계: 환자 ID(EMR 연동) 혼동 오류 방지, 직관적 '조기위험(Red)' 등급 가시화",
            "사이버보안: 환자 민감 정보 보호를 위한 데이터 종단간 암호화 및 EMR 망분리 준수 네트워크 구조 (KFDA 기준)",
            "로깅 및 이력 관리: 비인가자 열람 통제(RBAC 엑세스 제어), 열람 로그 100% 추적 아키텍처"
        ],
        "핵심 보안: 식약처 사이버보안 허가 가이드준수 및 의료진 피로도 저감 중심의 인간공학적 설계 탑재")

    # Slide 9: Roadmap
    create_content_slide(prs, 
        "08. 제품화 및 시장 진입(인허가, 사업화) 로드맵", 
        [
            "[2026.Q3] 제품 코어(엔진+UI) 동결 및 IEC 62304 기반 SW 밸리데이션(V&V) 완료",
            "[2026.Q4] GMP 및 SW 기술문서 심사 패키징, 식약처 품목 허가 심사 신청(혁신 콤보)",
            "[2027.Q1] 식약처 3등급 허가 획득 및 신의료기술 심사평가 진입 준비",
            "[2027.H2] 확증적 코호트 근거 추가 제출 및 선별 급여(혁신의료기술 지정) 추진"
        ],
        "Roadmap: 1년 내 인허가 패키징 문서화, 2년 내 임상 현장 비급여/혁신 급여 시장 진입 목표")

    # Slide 10: Conclusion
    create_content_slide(prs, 
        "09. 비전 종합 - 임상 현장에 정착될 ADDS V1.", 
        [
            "연구실 안의 AI를 넘어, 대장암 진료 현장에서 의사의 확신을 돕고 환자의 시간/비용을 아끼는 필수 진단 보조 장치",
            "환자 관점: 항암 치료 전 무의미한 항암 반복 배제 및 부작용 선제적 대비",
            "의료진 관점: 객관적 정량 지표(디지털 병리 + 영상 + 분자)에 기반한 강력한 의사결정 Reference",
            "산업 관점: 글로벌 수준(SaMD, GMP, Security)의 패키지 완성으로 디지털 치료/진단 시장 핵심 경쟁력 확보"
        ],
        "[최종 목표] 단순한 논문 실적(Publication)을 넘어 실제 환자를 살리고 규제를 돌파하는 지속 가능한 AI 의료기기 상용화")

    prs.save(out_path)
    print(f"Presentation saved successfully to {out_path}")

if __name__ == "__main__":
    out_dir = r"f:\ADDS\CDS"
    if not os.path.exists(out_dir):
        os.makedirs(out_dir)
        
    out_file = os.path.join(out_dir, "ADDS_SaMD_Pitch_Draft.pptx")
    build_presentation(out_file)
