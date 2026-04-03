from pptx import Presentation
import sys

def inspect_pptx(path):
    try:
        prs = Presentation(path)
        print(f"File: {path}")
        print(f"Number of slides: {len(prs.slides)}")
        for i, slide in enumerate(prs.slides):
            print(f"Slide {i+1}:")
            for shape in slide.shapes:
                print(f"  - Shape: {shape.shape_type}")
                if hasattr(shape, "text") and shape.text.strip():
                    print(f"    Text: {shape.text[:50]}...")
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    inspect_pptx("f:\\ADDS_CDS_PART2\\임상수요형중개연구 발표평가 준비.pptx")
